import io
from pathlib import Path
from pptx import Presentation


# -------------------------
#  DUPLICATE SLIDE FUNCTION
# -------------------------
def duplicate_slide(prs, index):
    """Duplicate an entire slide safely, including grouped shapes."""
    source = prs.slides[index]

    # Create a blank slide using the BLANK layout
    blank_slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Clear placeholder shapes
    blank_slide.shapes._spTree.clear()

    # Copy the full shape tree from the source slide
    for el in source.shapes._spTree:
        blank_slide.shapes._spTree.append(el)

    return blank_slide

# --------------------------------------------------
#  BUILD NBROWN GOVERNANCE PACK USING TEMPLATE SLIDES
# --------------------------------------------------
def build_governance_pack(
    client,
    period_start,
    period_end,
    sections,
    exec_summary,
    delivery_text,
    risks_issues_text,
    actions_text,
    total_weeks,
    current_week,
):
    """
    Build a Governance Pack using:
      Slide 1 → Title slide
      Slide 2 → Programme Status Report
      Slide 3 → Delivery Summary (extra workflows add Slides 4/5)
      Slide 4-7 → Project Delivery Overview (depends on workflows)
    """

    # -------------------------
    #  LOAD TEMPLATE
    # -------------------------
    modules_dir = Path(__file__).resolve().parent
    template_paths = [
        modules_dir / "clients" / "nbrown" / "templates" / "n_brown_status_update_template.pptx",
        modules_dir / "clients" / "nbrown" / "templates" / "gov_pack_template.pptx"
    ]

    template_path = next((p for p in template_paths if p.exists()), None)
    if not template_path:
        raise FileNotFoundError("NBrown PPT template not found.")

    prs = Presentation(str(template_path))

    # ----------------------------------------------------------
    #  SLIDE 1 — Title Slide (WE KEEP THE ORIGINAL)
    # ----------------------------------------------------------

    # ----------------------------------------------------------
    #  SLIDE 2 — Programme Status Report
    # ----------------------------------------------------------
    if sections.get("exec_summary"):
        slide = duplicate_slide(prs, 1)  # Slide 2 in PowerPoint is index 1

        print("\n\n--- SHAPES ON PROGRAMME STATUS SLIDE ---")
        for i, shape in enumerate(slide.shapes):
            text = shape.text if hasattr(shape, "text") else ""
            print(i, type(shape), text[:60])
        print("--- END SHAPES ---\n\n")

        # Find main body text box
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Executive Summary" in shape.text or "Status" in shape.text:
                    tf = shape.text_frame
                    tf.clear()
                    tf.text = exec_summary
                    break

    # ----------------------------------------------------------
    #  SLIDE 3 — Detailed Delivery Summary
    # ----------------------------------------------------------
    if sections.get("delivery_summary"):
        slide = duplicate_slide(prs, 2)  # Slide 3 = index 2
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Deliverables" in shape.text or "Summary" in shape.text:
                    tf = shape.text_frame
                    tf.clear()
                    tf.text = delivery_text
                    break

    # ----------------------------------------------------------
    #  SLIDE 4-7 — Timeline Slides (depending on workflows)
    # ----------------------------------------------------------
    if sections.get("timeline"):
        # Duplicate timeline slide once (you may duplicate more based on workflows)
        slide = duplicate_slide(prs, 3)  # Slide 4 = index 3
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "We are here" in shape.text:
                    tf = shape.text_frame
                    tf.clear()
                    tf.text = f"We are here → Week {current_week}"
                if "Timeline" in shape.text:
                    tf = shape.text_frame
                    tf.clear()
                    tf.text = f"Timeline {period_start} to {period_end}"

    # ----------------------------------------------------------
    #  SLIDE — Risks & Issues (generic NBrown layout)
    # ----------------------------------------------------------
    if sections.get("risks_issues"):
        slide = duplicate_slide(prs, 1)  # use generic NBrown content slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Weekly Status" in shape.text or "Status" in shape.text:
                    tf = shape.text_frame
                    tf.clear()
                    tf.text = risks_issues_text
                    break

    # ----------------------------------------------------------
    #  SLIDE — Actions Summary
    # ----------------------------------------------------------
    if sections.get("actions_summary"):
        slide = duplicate_slide(prs, 1)
        for shape in slide.shapes:
            if shape.has_text_frame:
                if "Weekly Status" in shape.text or "Status" in shape.text:
                    tf = shape.text_frame
                    tf.clear()
                    tf.text = actions_text
                    break

    # ----------------------------------------------------------
    # OUTPUT AS BYTES FOR STREAMLIT
    # ----------------------------------------------------------
    ppt_bytes = io.BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes
